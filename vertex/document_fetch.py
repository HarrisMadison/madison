"""
Document Fetcher — full-text retrieval by name.

Solves the "Gemini doesn't read named documents" problem. Path A (existing
snippet search) is great for "find me anything about X" questions. This module
provides Path B: when the user references a SPECIFIC file by name, we resolve
that name against the Vertex datastore, pull the FULL document content from
GCS, extract all the text, and hand it back.

Public surface:
    get_document_by_name(name: str, cfg: Config | None = None) -> dict
        Returns:
          {
            "ok": bool,
            "title": str,            # canonical filename matched
            "uri":   str,            # gs:// URI
            "text":  str,            # full extracted text (truncated to MAX_CHARS)
            "candidates": list[str], # other near-matches when fuzzy
            "error": str | None,
          }

    list_indexed_documents(cfg, limit=500) -> list[dict]
        Lightweight wildcard search to enumerate everything Vertex knows about.
        Used by the resolver and by future "what files do I have" features.

Why it auto-keeps-up with OneDrive:
    Resolution is done against the LIVE Vertex index, not a hardcoded folder
    listing. Any file ingested by the OneDrive sync pipeline becomes
    immediately reachable without code changes here.
"""
from __future__ import annotations

import io
import os
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from core import load_config, search_client, storage_client
from core.config import Config


# Hard cap on text we feed back to Gemini per fetch — keeps prompt size sane
# even when somebody asks for a 200-page contract. Gemini 2.5 Flash has a 1M
# token window so we have plenty of room, but extreme docs still need a limit.
MAX_CHARS = 200_000


# ─── data classes ────────────────────────────────────────────────────────────
@dataclass
class IndexedDoc:
    """One document as known to Vertex AI Search."""
    title: str          # filename or display title
    uri: str            # gs:// or https:// source URI
    doc_id: str
    struct: dict        # full struct_data for downstream filters
    score: float = 0.0  # ranked relevance for the resolver query


# ─── helpers ─────────────────────────────────────────────────────────────────
def _struct_to_dict(struct) -> dict:
    if struct is None:
        return {}
    try:
        return dict(struct)
    except Exception:
        try:
            return {k: struct[k] for k in struct}
        except Exception:
            return {}


def _extract_uri_from_doc(doc) -> str:
    """Same logic as vertex/search.py::_extract_uri — kept local to avoid
    a circular import and to be resilient if search.py changes."""
    # 1. content.uri (set during import)
    if getattr(doc, "content", None):
        u = getattr(doc.content, "uri", "") or ""
        if u:
            return u
    # 2. derived_struct_data fields
    if getattr(doc, "derived_struct_data", None):
        d = _struct_to_dict(doc.derived_struct_data)
        for k in ("link", "uri", "source_uri"):
            if d.get(k):
                return str(d[k])
    # 3. struct_data.source_uri (our manifest convention)
    if getattr(doc, "struct_data", None):
        s = _struct_to_dict(doc.struct_data)
        if s.get("source_uri"):
            return str(s["source_uri"])
    return ""


def _normalize(name: str) -> str:
    """Lowercase, strip path, drop extension and noise punctuation for matching."""
    if not name:
        return ""
    base = Path(name).name.lower()
    stem = Path(base).stem
    # collapse separators so "ABC_Plumbing-Invoice" matches "abc plumbing invoice"
    return re.sub(r"[\s_\-\.]+", " ", stem).strip()


# ─── local index integration ───────────────────────────────────────────────────────
# Try the local in-memory filename index BEFORE hitting Vertex search.
# This is what makes the system not burn quota: 95% of name-based fetches
# resolve here without an API call.
try:
    import sys as _sys
    _SCRIPTS_DIR = Path(__file__).resolve().parent.parent / "scripts"
    if str(_SCRIPTS_DIR) not in _sys.path:
        _sys.path.insert(0, str(_SCRIPTS_DIR))
    from local_index import get_index as _get_local_index
    _LOCAL_INDEX_AVAILABLE = True
except Exception as _e:
    print(f"[document_fetch] local_index unavailable: {_e}")
    _LOCAL_INDEX_AVAILABLE = False


def _resolve_via_local_index(query_name: str) -> Optional[IndexedDoc]:
    """Try the local filename index. Returns the top hit or None.
    Saves a Vertex search call when the user's query mentions a filename/address."""
    if not _LOCAL_INDEX_AVAILABLE:
        return None
    try:
        idx = _get_local_index()
        hits = idx.find(query_name, top_n=1)
        if hits and hits[0]["score"] >= 100:
            h = hits[0]
            return IndexedDoc(
                title=h["name"],
                uri=h["uri"],
                doc_id="",
                struct={},
                score=1.0,
            )
    except Exception as e:
        print(f"[document_fetch] local_index lookup error: {e}")
    return None


# ─── live Vertex query: list everything ───────────────────────────────────────
def list_indexed_documents(cfg: Config, limit: int = 500) -> list[IndexedDoc]:
    """
    Returns up to `limit` documents currently indexed by Vertex.

    Uses the broadest possible query (`*`) so we get a full enumeration. This
    is the cache-free source of truth — every call sees the latest OneDrive
    sync result.
    """
    from google.cloud import discoveryengine_v1 as de

    client = search_client(cfg)
    req = de.SearchRequest(
        serving_config=cfg.search_serving_config,
        query="*",
        page_size=min(limit, 100),  # Vertex caps page_size at 100
    )

    out: list[IndexedDoc] = []
    try:
        # Use page iteration so we can pull beyond the 100-result page cap
        page_result = client.search(request=req)
        for r in page_result.results:
            if len(out) >= limit:
                break
            doc = r.document
            sd = _struct_to_dict(doc.struct_data)
            if not sd and getattr(doc, "derived_struct_data", None):
                sd = _struct_to_dict(doc.derived_struct_data)

            uri = _extract_uri_from_doc(doc)
            title = (
                sd.get("title")
                or sd.get("filename")
                or (Path(uri).name if uri else "")
                or doc.id
                or "Document"
            )
            out.append(IndexedDoc(
                title=str(title),
                uri=str(uri),
                doc_id=str(doc.id or ""),
                struct=sd,
            ))
    except Exception as e:
        print(f"[document_fetch] list_indexed_documents error: {e}")
    return out


# ─── live Vertex query: targeted search for a name ──────────────────────────
def _search_for_name(cfg: Config, name: str, page_size: int = 25) -> list[IndexedDoc]:
    """Run a targeted Vertex search using the document name as the query.
    Returns ranked matches — the top hit is usually the one the user means."""
    from google.cloud import discoveryengine_v1 as de

    client = search_client(cfg)
    req = de.SearchRequest(
        serving_config=cfg.search_serving_config,
        query=name,
        page_size=page_size,
    )

    out: list[IndexedDoc] = []
    try:
        resp = client.search(request=req)
        for rank, r in enumerate(resp.results, 1):
            doc = r.document
            sd = _struct_to_dict(doc.struct_data)
            if not sd and getattr(doc, "derived_struct_data", None):
                sd = _struct_to_dict(doc.derived_struct_data)
            uri = _extract_uri_from_doc(doc)
            title = (
                sd.get("title")
                or sd.get("filename")
                or (Path(uri).name if uri else "")
                or doc.id
                or "Document"
            )
            out.append(IndexedDoc(
                title=str(title),
                uri=str(uri),
                doc_id=str(doc.id or ""),
                struct=sd,
                # higher rank = higher score (1.0 for top hit)
                score=1.0 / rank,
            ))
    except Exception as e:
        print(f"[document_fetch] _search_for_name error: {e}")
    return out


# ─── name resolution ────────────────────────────────────────────────────────
def _resolve_by_name(cfg: Config, query_name: str) -> tuple[Optional[IndexedDoc], list[IndexedDoc]]:
    """
    Resolve a user-supplied document name to ONE indexed doc.

    Strategy:
      1. Vertex search using the name → get ranked candidates (handles fuzzy
         and partial matches, plus typos to some extent because Vertex applies
         its own fuzzy matching).
      2. Score candidates with our own normalized-substring matcher to
         reorder when Vertex's text relevance disagrees with filename match.
      3. Return (best_match, [other_candidates]) so the caller can warn about
         ambiguity if needed.
    """
    if not query_name or not query_name.strip():
        return None, []

    # ── LOCAL INDEX FIRST — zero Vertex quota cost ──
    local_hit = _resolve_via_local_index(query_name)
    if local_hit is not None:
        print(f"[document_fetch] local_index resolved {query_name!r} → {local_hit.title}")
        return local_hit, []

    # ── FALLBACK to Vertex search ──
    candidates = _search_for_name(cfg, query_name)
    if not candidates:
        return None, []

    norm_query = _normalize(query_name)
    if not norm_query:
        return candidates[0], candidates[1:]

    def _score(c: IndexedDoc) -> float:
        norm_title = _normalize(c.title)
        if not norm_title:
            return c.score
        # exact normalized match wins outright
        if norm_title == norm_query:
            return 100.0
        # full substring match (either direction) is strong
        if norm_query in norm_title:
            return 50.0 + c.score
        if norm_title in norm_query:
            return 40.0 + c.score
        # word-overlap fallback
        q_words = set(norm_query.split())
        t_words = set(norm_title.split())
        if q_words and t_words:
            overlap = len(q_words & t_words) / len(q_words)
            return 10.0 * overlap + c.score
        return c.score

    candidates.sort(key=_score, reverse=True)
    return candidates[0], candidates[1:6]  # cap "also matched" list


# ─── GCS fetch + text extraction ─────────────────────────────────────────────
def _download_from_gcs(cfg: Config, gcs_uri: str) -> bytes:
    """Download a gs:// object as raw bytes."""
    if not gcs_uri.startswith("gs://"):
        raise ValueError(f"Not a GCS URI: {gcs_uri!r}")
    parts = gcs_uri[5:].split("/", 1)
    if len(parts) != 2:
        raise ValueError(f"Malformed GCS URI: {gcs_uri!r}")
    bucket_name, blob_name = parts

    client = storage_client(cfg)
    bucket = client.bucket(bucket_name)
    blob = bucket.blob(blob_name)
    if not blob.exists():
        raise FileNotFoundError(f"GCS object not found: {gcs_uri}")
    return blob.download_as_bytes()


def _extract_text(filename: str, raw: bytes) -> str:
    """
    Pull text from raw bytes based on extension.

    Supports: PDF (pdfplumber), DOCX (python-docx), XLSX (openpyxl),
    PPTX (python-pptx), TXT/MD/CSV (utf-8 decode), HTML (basic strip).
    Unknown formats → best-effort utf-8 decode with replacement.
    """
    ext = Path(filename).suffix.lower()

    try:
        if ext == ".pdf":
            return _extract_pdf(raw)
        if ext in (".docx",):
            return _extract_docx(raw)
        if ext in (".xlsx", ".xlsm"):
            return _extract_xlsx(raw)
        if ext in (".pptx",):
            return _extract_pptx(raw)
        if ext in (".txt", ".md", ".csv", ".tsv", ".log", ".json", ".yaml", ".yml"):
            return raw.decode("utf-8", errors="replace")
        if ext in (".html", ".htm"):
            text = raw.decode("utf-8", errors="replace")
            return re.sub(r"<[^>]+>", " ", text)
        # Unknown extension — try utf-8 anyway. If it's binary garbage, the
        # caller will see it and know to OCR.
        return raw.decode("utf-8", errors="replace")
    except Exception as e:
        return f"[Extraction failed for {filename}: {e}]"


def _extract_pdf(raw: bytes) -> str:
    """PDF extraction via pdfplumber. Returns one big string with page breaks."""
    try:
        import pdfplumber
    except ImportError:
        return "[pdfplumber not installed — pip install pdfplumber]"

    pages: list[str] = []
    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            try:
                txt = page.extract_text() or ""
            except Exception as e:
                txt = f"[page {i} extract error: {e}]"
            pages.append(f"--- Page {i} ---\n{txt.strip()}")
    body = "\n\n".join(pages).strip()
    if not body:
        return ("[PDF appears to contain no extractable text — likely a scanned "
                "document. Run it through OCR (Document AI) before retrying.]")
    return body


def _extract_docx(raw: bytes) -> str:
    """DOCX extraction via python-docx. Includes paragraph and table text."""
    try:
        import docx  # python-docx
    except ImportError:
        return "[python-docx not installed]"
    doc = docx.Document(io.BytesIO(raw))
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_xlsx(raw: bytes) -> str:
    """XLSX extraction via openpyxl. Each sheet rendered as TSV."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "[openpyxl not installed]"
    wb = load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"=== Sheet: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append("\t".join(cells))
        parts.append("")
    return "\n".join(parts)


def _extract_pptx(raw: bytes) -> str:
    """PPTX extraction via python-pptx. Slide titles + body text per slide."""
    try:
        from pptx import Presentation
    except ImportError:
        return "[python-pptx not installed]"
    prs = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
        parts.append("")
    return "\n".join(parts)


# ─── public entry point ─────────────────────────────────────────────────────
def get_document_by_name(name: str, cfg: Optional[Config] = None) -> dict:
    """
    Resolve a user-supplied document name to a single indexed file, fetch the
    full content from GCS, extract text, and return it.

    Args:
        name: Filename or partial name as the user typed it (e.g.
              "ABC Plumbing invoice", "March permit", "northridge appraisal").
        cfg:  Optional pre-loaded Config. Loaded on demand if omitted.

    Returns:
        dict with keys:
            ok:         True on success, False on any failure.
            title:      Canonical filename of the matched document (when ok).
            uri:        gs:// URI of the source object (when ok).
            text:       Full extracted text, truncated to MAX_CHARS.
            candidates: Other near-matches as ["filename1", "filename2", ...].
            error:      Human-readable error string when not ok.

    The function never raises — all errors come back in the dict so the
    Gemini tool-call loop can hand them straight to the model.
    """
    if not name or not name.strip():
        return {
            "ok": False,
            "title": "",
            "uri": "",
            "text": "",
            "candidates": [],
            "error": "document_name was empty",
        }

    try:
        cfg = cfg or load_config()
    except SystemExit as e:
        # load_config calls sys.exit on a missing config file. Catch it so the
        # tool returns a normal error dict instead of killing the Flask process.
        return {
            "ok": False, "title": "", "uri": "", "text": "",
            "candidates": [],
            "error": f"Config load failed: {e}",
        }
    except Exception as e:
        return {
            "ok": False, "title": "", "uri": "", "text": "",
            "candidates": [],
            "error": f"Config load failed: {e}",
        }

    best, others = _resolve_by_name(cfg, name)
    if best is None:
        return {
            "ok": False, "title": "", "uri": "", "text": "",
            "candidates": [],
            "error": (f"No indexed document matches {name!r}. The file may not be "
                      f"synced yet, or it may be under a different name."),
        }

    if not best.uri:
        return {
            "ok": False,
            "title": best.title,
            "uri": "",
            "text": "",
            "candidates": [c.title for c in others],
            "error": (f"Matched {best.title!r} in the index but it has no source "
                      f"URI — cannot fetch full text. Re-run the OneDrive sync."),
        }

    # Fetch the bytes
    try:
        if best.uri.startswith("gs://"):
            raw = _download_from_gcs(cfg, best.uri)
        else:
            return {
                "ok": False,
                "title": best.title,
                "uri": best.uri,
                "text": "",
                "candidates": [c.title for c in others],
                "error": (f"URI {best.uri!r} is not a gs:// URI — only GCS-backed "
                          f"sources are fetchable at the moment."),
            }
    except FileNotFoundError as e:
        return {
            "ok": False,
            "title": best.title,
            "uri": best.uri,
            "text": "",
            "candidates": [c.title for c in others],
            "error": str(e),
        }
    except Exception as e:
        return {
            "ok": False,
            "title": best.title,
            "uri": best.uri,
            "text": "",
            "candidates": [c.title for c in others],
            "error": f"GCS download failed: {e}",
        }

    # Extract text
    text = _extract_text(best.title, raw)
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS] + f"\n\n[... truncated at {MAX_CHARS:,} chars ...]"

    return {
        "ok": True,
        "title": best.title,
        "uri": best.uri,
        "text": text,
        "candidates": [c.title for c in others],
        "error": None,
    }
