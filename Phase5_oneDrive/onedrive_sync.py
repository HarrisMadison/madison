#!/usr/bin/env python3
"""
Phase5_OneDrive/onedrive_sync.py
----------------------------------
Syncs files from a OneDrive folder to a GCS bucket, then triggers
a Vertex AI Search re-import.

Usage:
  python onedrive_sync.py               # incremental sync (delta)
  python onedrive_sync.py --force       # full re-sync regardless of delta
  python onedrive_sync.py --dry-run     # list what would sync, no writes
  python onedrive_sync.py --schedule 30 # loop every 30 minutes

SCALE-TODO: switch to client_credentials before production.
See bootstrap_onedrive.py for full instructions.
"""

import os, sys, json, time, logging, argparse, requests, msal, base64, io
# Phase 6: lazy-loaded OCR + metadata enrichment
_P6_LOADED = False
_enrich_metadata = None
_needs_ocr       = None
_ocr_pdf_gcs     = None

def _load_phase6():
    global _P6_LOADED, _enrich_metadata, _needs_ocr, _ocr_pdf_gcs
    if _P6_LOADED:
        return
    try:
        from phase6_ocr_metadata import enrich_metadata, needs_ocr, ocr_pdf_gcs
        _enrich_metadata = enrich_metadata
        _needs_ocr       = needs_ocr
        _ocr_pdf_gcs     = ocr_pdf_gcs
        log.info("Phase 6 OCR+metadata loaded.")
    except ImportError:
        pass
    _P6_LOADED = True
from datetime import datetime, timezone
from pathlib import Path
from google.cloud import storage
import google.auth
import google.auth.transport.requests

def _load_env():
    from dotenv import load_dotenv
    candidates = [
        os.environ.get("VERTEX_ENV_FILE"),
        Path(__file__).parent / "secrets" / ".env",
        Path(__file__).parent / "Secrets" / ".env",   # Windows-cased variant
        Path.cwd() / "Phase5_oneDrive" / "secrets" / ".env",
        Path.cwd() / ".env",
        # Repo top-level .env (this is where Madison Ave's real env actually lives)
        Path(__file__).resolve().parent.parent / ".env",
        Path(__file__).resolve().parent.parent / "Phase3_Bootstrap" / "secrets" / ".env",
    ]
    for c in candidates:
        if c and Path(c).exists():
            load_dotenv(c)
            print(f"  Loaded env: {c}")
            return
    print("  WARNING: no .env file found — AZURE_TENANT_ID etc. will be empty")

_load_env()

AZURE_CLIENT_ID      = os.environ.get("AZURE_CLIENT_ID", "")
AZURE_TENANT_ID      = os.environ.get("AZURE_TENANT_ID", "")
ONEDRIVE_FOLDER_PATH = os.environ.get("ONEDRIVE_FOLDER_PATH", "")
GCP_PROJECT_ID       = os.environ.get("GCP_PROJECT_ID", "")
GCS_BUCKET_NAME      = os.environ.get("GCS_BUCKET_NAME", "")
VERTEX_LOCATION      = os.environ.get("VERTEX_LOCATION", "global")
VERTEX_DATASTORE     = os.environ.get("VERTEX_DATASTORE_ID", "")

SCOPES           = ["Files.Read"]
TOKEN_CACHE_PATH = Path(__file__).parent / "secrets" / "token_cache.json"
DELTA_STATE_PATH = Path(__file__).parent / "secrets" / "delta_state.json"
GRAPH_API        = "https://graph.microsoft.com/v1.0"

logging.basicConfig(level=logging.INFO, format="%(asctime)s  %(levelname)-8s  %(message)s", datefmt="%Y-%m-%d %H:%M:%S")
log = logging.getLogger("onedrive_sync")

def _get_ms_token():
    cache = msal.SerializableTokenCache()
    if TOKEN_CACHE_PATH.exists():
        cache.deserialize(TOKEN_CACHE_PATH.read_text())
    app = msal.PublicClientApplication(
        AZURE_CLIENT_ID,
        authority=f"https://login.microsoftonline.com/{AZURE_TENANT_ID}",
        token_cache=cache,
    )
    accounts = app.get_accounts()
    result = None
    if accounts:
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result:
            log.info("Auth: using cached token")
    if not result:
        log.warning("Cached token unavailable -- starting device-code flow")
        flow = app.initiate_device_flow(scopes=SCOPES)
        if "user_code" not in flow:
            raise RuntimeError(f"Device flow failed: {flow.get('error_description')}")
        print("\n" + "="*60)
        print("  ACTION REQUIRED -- Microsoft sign-in")
        print("="*60)
        print(f"  1. Open:       {flow['verification_uri']}")
        print(f"  2. Enter code: {flow['user_code']}")
        print("  3. Sign in with the OneDrive account")
        print("="*60 + "\n")
        result = app.acquire_token_by_device_flow(flow)
    if "access_token" not in result:
        raise RuntimeError(
            f"Microsoft auth failed: {result.get('error_description', result)}\n"
            "SCALE-TODO: Token may have expired. See bootstrap_onedrive.py."
        )
    TOKEN_CACHE_PATH.parent.mkdir(exist_ok=True)
    TOKEN_CACHE_PATH.write_text(cache.serialize())
    return result["access_token"]

def _get_drive_id(token):
    """Get OneDrive for Business drive ID -- required for SharePoint-hosted drives."""
    r = requests.get(f"{GRAPH_API}/me/drive", headers={"Authorization": f"Bearer {token}"})
    r.raise_for_status()
    return r.json()["id"]

def _get_gcp_token():
    creds, _ = google.auth.default(scopes=["https://www.googleapis.com/auth/cloud-platform"])
    creds.refresh(google.auth.transport.requests.Request())
    return creds.token

def list_onedrive_files(token, force):
    delta_state = {}
    if DELTA_STATE_PATH.exists() and not force:
        try:
            delta_state = json.loads(DELTA_STATE_PATH.read_text())
        except Exception:
            pass

    delta_link = delta_state.get("delta_link")
    files = []

    if delta_link and not force:
        log.info("Using OneDrive delta link (incremental sync)")
        url = delta_link
    else:
        scope_label = ONEDRIVE_FOLDER_PATH if ONEDRIVE_FOLDER_PATH else "(entire OneDrive root)"
        log.info(f"Full listing scope: {scope_label}")
        drive_id = _get_drive_id(token)
        log.info(f"Drive ID: {drive_id[:20]}...")
        # Empty ONEDRIVE_FOLDER_PATH -> sync the whole drive from root.
        # Non-empty -> scope to that folder path (legacy DoorLoop behavior).
        if ONEDRIVE_FOLDER_PATH.strip():
            url = f"{GRAPH_API}/drives/{drive_id}/root:/{ONEDRIVE_FOLDER_PATH}:/delta"
        else:
            url = f"{GRAPH_API}/drives/{drive_id}/root/delta"

    while url:
        # Retry loop with backoff for Graph API rate limiting (429)
        for attempt in range(5):
            r = requests.get(url, headers={"Authorization": f"Bearer {token}"})
            if r.status_code == 429:
                wait = int(r.headers.get("Retry-After", 20)) + 2
                log.warning(f"Rate limited by Graph API -- waiting {wait}s (attempt {attempt+1}/5)")
                time.sleep(wait)
                continue
            r.raise_for_status()
            break
        data = r.json()
        for item in data.get("value", []):
            if "deleted" in item:
                log.info(f"  Deleted on OneDrive (skipping): {item.get('name', item['id'])}")
                continue
            if "file" in item:
                files.append(item)
        url = data.get("@odata.nextLink")
        new_delta = data.get("@odata.deltaLink")
        if new_delta:
            delta_state["delta_link"] = new_delta

    DELTA_STATE_PATH.parent.mkdir(exist_ok=True)
    DELTA_STATE_PATH.write_text(json.dumps(delta_state, indent=2))
    log.info(f"OneDrive: {len(files)} file(s) to sync")
    return files

def download_file(token, item):
    download_url = item.get("@microsoft.graph.downloadUrl")
    if not download_url:
        download_url = f"{GRAPH_API}/me/drive/items/{item['id']}/content"
    r = requests.get(download_url, headers={"Authorization": f"Bearer {token}"}, stream=True)
    # If 401, silently refresh token and retry once
    if r.status_code == 401:
        log.warning("Token expired mid-sync -- refreshing token and retrying")
        token = _get_ms_token()
        r = requests.get(download_url, headers={"Authorization": f"Bearer {token}"}, stream=True)
    r.raise_for_status()
    return r.content, token

def upload_to_gcs(data, filename, item, dry_run):
    # Photos are tracked via pointer docs only -- skip GCS upload entirely
    if any(filename.lower().endswith(ext) for ext in _PHOTO_EXTS):
        return None
    # Preserve OneDrive path structure using parentReference
    parent_path = item.get("parentReference", {}).get("path", "").split("root:")[-1].strip("/")
    gcs_path = f"onedrive-mirror/{parent_path}/{filename}" if parent_path else f"onedrive-mirror/{filename}"
    uri = f"gs://{GCS_BUCKET_NAME}/{gcs_path}"
    if dry_run:
        log.info(f"  [dry-run] Would upload -> {uri}")
        return uri
    client = storage.Client(project=GCP_PROJECT_ID)
    bucket = client.bucket(GCS_BUCKET_NAME)
    bucket.blob(gcs_path).upload_from_string(data)
    log.info(f"  Uploaded -> {uri}")
    return uri

# Searchable document types for Vertex AI Search
_SEARCHABLE_EXTS = ('.pdf', '.docx', '.doc', '.xlsx', '.xls', '.csv', '.txt', '.pptx')
_MIME_MAP = {
    '.pdf':  'application/pdf',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.doc':  'application/msword',
    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    '.xls':  'application/vnd.ms-excel',
    '.csv':  'text/csv',
    '.txt':  'text/plain',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
}

# Photo extensions -- skipped from GCS upload; pointer docs created instead
_PHOTO_EXTS = ('.jpg', '.jpeg', '.png', '.heic', '.heif', '.tiff', '.tif', '.webp')

# PDFs larger than this become pointer-only docs in the manifest (no content extraction)
_LARGE_PDF_BYTES = 8 * 1024 * 1024   # 8 MB

# Module-level state: set by run_sync, consumed by trigger_vertex_import
_last_token    = ""
_last_drive_id = ""
_last_items: list = []


# ── Pre-extraction support ───────────────────────────────────────────────────────
# Pulls text out of PDF/DOCX/XLSX/CSV/TXT directly so we hand Vertex
# already-extracted text via rawBytes instead of relying on its own parser.
# This solves the core problem where Vertex's parser silently fails on
# certain text PDFs (invoicing software output, hybrid PDFs, etc.) and
# returns empty snippets even though the file has perfectly readable text.
#
# Falls back to letting Vertex parse the file when our extractor returns
# nothing -- so this is purely additive, never destructive.

MAX_EXTRACTED_CHARS = 200_000   # cap per doc to avoid huge proto payloads

def _extract_text_from_bytes(file_bytes: bytes, ext: str, name_for_log: str = "") -> str:
    """Extract plain text from common doc formats. Returns '' on failure."""
    ext = ext.lower()
    try:
        if ext == ".pdf":
            try:
                import pdfplumber
                text_parts = []
                with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                    for page in pdf.pages:
                        t = page.extract_text() or ""
                        if t:
                            text_parts.append(t)
                        # Also try to capture tables as readable text
                        for table in page.extract_tables() or []:
                            for row in table:
                                cells = [c for c in (row or []) if c]
                                if cells:
                                    text_parts.append(" | ".join(str(c) for c in cells))
                return "\n".join(text_parts).strip()
            except Exception as e:
                log.debug(f"  pdfplumber failed on {name_for_log}: {e}")
                return ""

        if ext in (".docx",):
            try:
                from docx import Document  # python-docx
                doc = Document(io.BytesIO(file_bytes))
                parts = [p.text for p in doc.paragraphs if p.text]
                # Tables
                for table in doc.tables:
                    for row in table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                return "\n".join(parts).strip()
            except ImportError:
                log.debug("  python-docx not installed; .docx will use Vertex parser")
                return ""
            except Exception as e:
                log.debug(f"  docx parse failed on {name_for_log}: {e}")
                return ""

        if ext in (".xlsx", ".xls"):
            try:
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
                parts = []
                for sheet in wb.worksheets:
                    parts.append(f"[Sheet: {sheet.title}]")
                    for row in sheet.iter_rows(values_only=True):
                        cells = [str(c) for c in row if c is not None and str(c).strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                return "\n".join(parts).strip()
            except ImportError:
                log.debug("  openpyxl not installed; .xlsx will use Vertex parser")
                return ""
            except Exception as e:
                log.debug(f"  xlsx parse failed on {name_for_log}: {e}")
                return ""

        if ext in (".csv", ".txt"):
            try:
                return file_bytes.decode("utf-8", errors="replace").strip()
            except Exception:
                return ""

        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                parts = []
                for i, slide in enumerate(prs.slides, 1):
                    parts.append(f"[Slide {i}]")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            parts.append(shape.text)
                return "\n".join(parts).strip()
            except ImportError:
                log.debug("  python-pptx not installed; .pptx will use Vertex parser")
                return ""
            except Exception as e:
                log.debug(f"  pptx parse failed on {name_for_log}: {e}")
                return ""

    except Exception as e:
        log.debug(f"  Unexpected extraction error on {name_for_log}: {e}")
    return ""


def _build_item_index(items: list) -> dict:
    """Map GCS blob path -> { webUrl, item_id, name } from OneDrive items.
    Used at manifest-build time to inject OneDrive 'open in browser' URLs into
    every doc's structData (the DoorLoop pipeline does the equivalent with a
    Drive cache).
    """
    out = {}
    for it in items:
        if "file" not in it:
            continue
        name = it.get("name", "")
        parent_path = it.get("parentReference", {}).get("path", "").split("root:")[-1].strip("/")
        gcs_key = f"onedrive-mirror/{parent_path}/{name}" if parent_path else f"onedrive-mirror/{name}"
        out[gcs_key] = {
            "webUrl":  it.get("webUrl", ""),
            "item_id": it.get("id", ""),
            "name":    name,
        }
    return out


def _make_doc_id(blob_name: str) -> str:
    """Sanitize a GCS blob path into a valid Vertex document ID.
    Vertex requires: [a-zA-Z0-9-_]* only.
    """
    import re
    clean = re.sub(r'[^a-zA-Z0-9_]', '_', blob_name)
    clean = re.sub(r'_+', '_', clean).strip('_')
    return clean[:128]


def _get_onedrive_folder_url(token: str, drive_id: str, folder_path: str) -> str:
    """Get the webUrl for a OneDrive folder path. Returns empty string on failure."""
    try:
        url = f"{GRAPH_API}/drives/{drive_id}/root:/{folder_path}"
        r = requests.get(url, headers={"Authorization": f"Bearer {token}"})
        if r.ok:
            return r.json().get("webUrl", "")
    except Exception:
        pass
    return ""


def _build_photo_pointer_docs(token: str, drive_id: str, folder_path: str, items: list) -> list:
    """
    Group photo items by property folder and build one Vertex pointer doc per property.
    Each pointer doc contains the photo count and a direct OneDrive URL so Bob
    can navigate straight to the photos from a chat answer.
    """
    from collections import defaultdict
    property_photos: dict = defaultdict(list)

    for item in items:
        name = item.get("name", "")
        if not any(name.lower().endswith(ext) for ext in _PHOTO_EXTS):
            continue
        parent_path = item.get("parentReference", {}).get("path", "").split("root:")[-1].strip("/")
        parts = [p for p in parent_path.split("/") if p]

        # Decide grouping key:
        # - If ONEDRIVE_FOLDER_PATH is set (e.g. "Doorloop"), strip it from the
        #   front of the path and use the next 1-2 segments as the property key.
        #   This preserves the original DoorLoop layout behavior.
        # - If ONEDRIVE_FOLDER_PATH is empty (whole-drive sync), group by the
        #   immediate parent folder of the photo (or up to 2 trailing segments).
        if folder_path.strip():
            scope_parts = [p for p in folder_path.split("/") if p]
            # Drop the scope prefix from parts if present
            if parts[: len(scope_parts)] == scope_parts:
                rel = parts[len(scope_parts):]
            else:
                rel = parts
            if len(rel) >= 2:
                prop_key = f"{rel[0]}/{rel[1]}"
            elif len(rel) == 1:
                prop_key = rel[0]
            else:
                prop_key = "Unknown"
        else:
            # Whole-drive mode: group by immediate parent folder; if there is a
            # grandparent we keep both so we can reconstruct an OneDrive URL.
            if len(parts) >= 2:
                prop_key = f"{parts[-2]}/{parts[-1]}"
            elif len(parts) == 1:
                prop_key = parts[0]
            else:
                prop_key = "Unknown"

        property_photos[prop_key].append(item)

    pointer_docs = []
    for prop_key, photo_items in sorted(property_photos.items()):
        parts       = prop_key.split("/")
        prop_name   = parts[0]
        sub_folder  = parts[1] if len(parts) > 1 else "photos"
        photo_count = len(photo_items)

        od_folder_path = "/".join([s for s in [folder_path, prop_name, sub_folder] if s])
        od_url = _get_onedrive_folder_url(token, drive_id, od_folder_path)

        doc_id    = _make_doc_id(f"photo_pointer_{prop_key}")
        json_data = {
            "title":          f"{prop_name} — Photos ({photo_count} images)",
            "property":       prop_name,
            "document_type":  "photo_index",
            "photo_count":    photo_count,
            "onedrive_url":   od_url,
            "summary": (
                f"There are {photo_count} photos for {prop_name} stored in OneDrive. "
                f"Click here to view them: {od_url}"
            ),
        }
        # Write BOTH structData (for Vertex ranking/filtering) and
        # jsonData (for full-text search content). structData lets
        # retrieve() detect photo_index docs by document_type field.
        struct_data = {
            "title":         json_data["title"],
            "property":      prop_name,
            "document_type": "photo_index",
            "photo_count":   photo_count,
            "onedrive_url":  od_url,
            "source_uri":    od_url,
        }
        pointer_docs.append({
            "id":         doc_id,
            "jsonData":   json.dumps(json_data),
            "structData": json.dumps(struct_data),
        })
        log.info(f"  Photo pointer: {prop_name} ({photo_count} photos)")

    return pointer_docs


def _build_and_upload_manifest(dry_run: bool, token: str = "", drive_id: str = "", items: list = None) -> str | None:
    """
    Build a Vertex import manifest with three document types:
      1. Regular docs (PDF/DOCX/XLSX etc) -- full content extraction
      2. Large PDF pointers (>8MB) -- metadata + GCS link only
      3. Photo pointer docs -- one per property with OneDrive URL
    """
    if dry_run:
        log.info("  [dry-run] Would build and upload Vertex import manifest")
        return None

    import google.auth
    creds, _ = google.auth.default(scopes=["https://www.googleapis.com/auth/cloud-platform"])
    gcs_client = storage.Client(project=GCP_PROJECT_ID, credentials=creds)
    bucket     = gcs_client.bucket(GCS_BUCKET_NAME)

    # Build OneDrive URL index so we can stamp webUrl into every doc
    item_index = _build_item_index(items or [])
    log.info(f"  Item index built: {len(item_index)} OneDrive URLs available")

    lines: list[str] = []
    seen_ids: set    = set()
    count_docs = count_large = 0
    count_extracted = count_ocr = count_passthrough = 0

    for blob in bucket.list_blobs(prefix="onedrive-mirror/"):
        name_lower = blob.name.lower()
        ext = next((e for e in _SEARCHABLE_EXTS if name_lower.endswith(e)), None)
        if not ext:
            continue
        uri    = f"gs://{GCS_BUCKET_NAME}/{blob.name}"
        title  = blob.name.split("/")[-1]
        doc_id = _make_doc_id(blob.name)
        if doc_id in seen_ids:
            doc_id = f"{doc_id[:120]}_{len(seen_ids)}"
        seen_ids.add(doc_id)

        # Large PDFs: pointer doc only (no content extraction)
        if ext == ".pdf" and blob.size and blob.size > _LARGE_PDF_BYTES:
            size_mb   = blob.size / (1024 * 1024)
            base_meta = {
                "title":         title,
                "document_type": "large_pdf_pointer",
                "size_mb":       round(size_mb, 1),
                "gcs_uri":       uri,
                "summary":       f"{title} is a {size_mb:.1f} MB PDF. GCS path: {uri}",
            }
            # Stamp OneDrive web URL if available
            od_info = item_index.get(blob.name)
            if od_info and od_info.get("webUrl"):
                base_meta["onedrive_url"] = od_info["webUrl"]
                base_meta["source_uri"]   = od_info["webUrl"]
            _load_phase6()
            if _enrich_metadata:
                base_meta = _enrich_metadata(blob.name, base_meta)
                base_meta["document_type"] = "large_pdf_pointer"
            lines.append(json.dumps({"id": doc_id, "jsonData": json.dumps(base_meta)}))
            count_large += 1
        else:
            base_struct = {"title": title, "source_uri": uri}
            # Stamp OneDrive web URL if available so the chat can link directly
            od_info = item_index.get(blob.name)
            if od_info and od_info.get("webUrl"):
                base_struct["onedrive_url"] = od_info["webUrl"]
                # Prefer OneDrive URL as source_uri (better than gs:// for users)
                base_struct["source_uri"]   = od_info["webUrl"]
                base_struct["gcs_uri"]      = uri

            _load_phase6()
            if _enrich_metadata:
                base_struct = _enrich_metadata(blob.name, base_struct)

            # ── Step 1: try our own pre-extraction (pdfplumber / docx / xlsx) ────
            extracted_text = ""
            try:
                file_bytes = blob.download_as_bytes()
                extracted_text = _extract_text_from_bytes(file_bytes, ext, blob.name)
            except Exception as e:
                log.debug(f"  Could not download {blob.name} for extraction: {e}")

            # ── Step 2: if extraction failed and OCR is enabled, try OCR ────────
            ocr_text = None
            if not extracted_text and _needs_ocr and _ocr_pdf_gcs and ext == ".pdf":
                if _needs_ocr(blob.name, blob.size or 0):
                    # Pass bucket + source metadata so the OCR call can read/write
                    # the gs://<bucket>/ocr-cache/ entry. This makes each scan OCR'd
                    # exactly ONCE -- subsequent syncs read from cache for free.
                    ocr_text = _ocr_pdf_gcs(
                        uri, GCP_PROJECT_ID,
                        bucket=bucket,
                        source_blob_name=blob.name,
                        source_updated=blob.updated,
                    )

            # ── Decide which content path to ship to Vertex ────────────────────
            content_text = extracted_text or ocr_text or ""
            if content_text:
                if len(content_text) > MAX_EXTRACTED_CHARS:
                    content_text = content_text[:MAX_EXTRACTED_CHARS]
                lines.append(json.dumps({
                    "id":       doc_id,
                    "jsonData": json.dumps(base_struct),
                    "content":  {
                        "mimeType": "text/plain",
                        "rawBytes": base64.b64encode(
                            content_text.encode("utf-8")
                        ).decode("ascii"),
                    },
                }))
                if extracted_text:
                    count_extracted += 1
                else:
                    count_ocr += 1
            else:
                # Fall through: let Vertex's own parser try the file.
                lines.append(json.dumps({
                    "id":       doc_id,
                    "jsonData": json.dumps(base_struct),
                    "content":  {"mimeType": _MIME_MAP.get(ext, "application/pdf"), "uri": uri},
                }))
                count_passthrough += 1
            count_docs += 1

    # Photo pointer docs (one per property, with OneDrive URL)
    photo_pointers = []
    if token and drive_id and items:
        photo_pointers = _build_photo_pointer_docs(token, drive_id, ONEDRIVE_FOLDER_PATH, items)
        for p in photo_pointers:
            if p["id"] not in seen_ids:
                seen_ids.add(p["id"])
                lines.append(json.dumps(p))
    count_photos = len(photo_pointers)

    manifest_path = "manifests/import_manifest_latest.jsonl"
    bucket.blob(manifest_path).upload_from_string("\n".join(lines))
    manifest_uri  = f"gs://{GCS_BUCKET_NAME}/{manifest_path}"

    # Write photo_index.json — flat lookup for Bob to read directly from GCS.
    # Keyed by property name, value has OneDrive URL + photo count.
    # This avoids Vertex search quota for photo lookups entirely.
    if photo_pointers:
        photo_index = {}
        for p in photo_pointers:
            try:
                data = json.loads(p.get("jsonData", "{}"))
                prop = data.get("property", "")
                url  = data.get("onedrive_url", "")
                cnt  = data.get("photo_count", 0)
                if prop and url:
                    photo_index[prop] = {"url": url, "count": cnt, "title": data.get("title", "")}
            except Exception:
                pass
        if photo_index:
            bucket.blob("manifests/photo_index.json").upload_from_string(
                json.dumps(photo_index, indent=2))
            log.info(f"  Photo index written: {len(photo_index)} properties -> "
                     f"gs://{GCS_BUCKET_NAME}/manifests/photo_index.json")

    log.info(
        f"Manifest: {count_docs} docs ({count_extracted} pre-extracted, "
        f"{count_ocr} OCR, {count_passthrough} passthrough) + "
        f"{count_large} large-PDF pointers + {count_photos} photo pointers "
        f"-> {manifest_uri}"
    )
    return manifest_uri


def trigger_vertex_import(dry_run):
    if not VERTEX_DATASTORE or not GCP_PROJECT_ID:
        log.warning("VERTEX_DATASTORE_ID or GCP_PROJECT_ID not set -- skipping Vertex import")
        return

    manifest_uri = _build_and_upload_manifest(dry_run, token=_last_token, drive_id=_last_drive_id, items=_last_items)
    if dry_run:
        return

    url = (
        f"https://discoveryengine.googleapis.com/v1alpha/projects/{GCP_PROJECT_ID}"
        f"/locations/{VERTEX_LOCATION}/collections/default_collection"
        f"/dataStores/{VERTEX_DATASTORE}/branches/0/documents:import"
    )
    body = {
        "gcsSource": {
            "inputUris": [manifest_uri],
            "dataSchema": "document",
        },
        "reconciliationMode": "FULL",
    }
    token = _get_gcp_token()
    headers = {
        "Authorization": f"Bearer {token}",
        "X-Goog-User-Project": GCP_PROJECT_ID,
        "Content-Type": "application/json",
    }
    r = requests.post(url, headers=headers, json=body)
    if r.status_code == 200:
        log.info(f"Vertex import triggered. Operation: {r.json().get('name', '')}")
    else:
        log.error(f"Vertex import failed: {r.status_code} {r.text}")

def run_sync(dry_run=False, force=False, rebuild_only=False):
    global _last_token, _last_drive_id, _last_items

    log.info("=" * 50)
    log.info(f"OneDrive sync started -- dry_run={dry_run}, force={force}, rebuild_only={rebuild_only}")
    log.info("=" * 50)

    ms_token = _get_ms_token()
    _last_token = ms_token

    try:
        _last_drive_id = _get_drive_id(ms_token)
    except Exception:
        _last_drive_id = ""

    if rebuild_only:
        # Skip the OneDrive download entirely. GCS already has every file
        # mirrored. We just need to walk GCS, run pre-extraction + OCR,
        # build the manifest, and trigger Vertex re-import.
        # We still pull OneDrive item metadata so we can stamp webUrls and
        # build photo pointer docs, but no file bytes get downloaded from
        # OneDrive.
        log.info("Rebuild-only mode: skipping OneDrive file download.")
        log.info("Fetching OneDrive item index for webUrl stamping...")
        files = list_onedrive_files(ms_token, force=True)
        _last_items = files
        log.info(f"OneDrive metadata loaded: {len(files)} items (used for webUrl + photo pointers only)")
        log.info("Building manifest from existing GCS contents with pre-extraction + OCR...")
        trigger_vertex_import(dry_run=dry_run)
        log.info("=" * 50)
        return

    files = list_onedrive_files(ms_token, force=force)
    _last_items = files

    if not files:
        log.info("No files to sync.")
        return

    photos = [f for f in files if any(f["name"].lower().endswith(e) for e in _PHOTO_EXTS)]
    docs   = [f for f in files if not any(f["name"].lower().endswith(e) for e in _PHOTO_EXTS)]
    log.info(f"Plan: {len(docs)} documents to upload + {len(photos)} photos (skipped -- pointer docs only)")

    uploaded = errors = 0
    for item in docs:
        name    = item["name"]
        size_kb = item.get("size", 0) // 1024
        log.info(f"Syncing: {name}  ({size_kb} KB)")
        try:
            if not dry_run:
                data, ms_token = download_file(ms_token, item)
                _last_token = ms_token
                upload_to_gcs(data, name, item, dry_run=False)
            else:
                log.info(f"  [dry-run] Would upload -> {name}")
            uploaded += 1
        except Exception as e:
            log.error(f"  Failed: {name} -- {e}")
            errors += 1

    log.info(f"Sync complete: {uploaded} docs uploaded, {len(photos)} photos skipped, {errors} errors")
    # Always trigger Vertex import (even when --force has no new files) so the
    # manifest is rebuilt with the latest pre-extraction + OCR logic.
    if uploaded > 0 or dry_run or force:
        trigger_vertex_import(dry_run=dry_run)
    log.info("=" * 50)

def trigger_local_index_reload():
    """Tell the running Flask server (simple_web.py) to reload its in-memory
    filename index, so newly synced files become searchable in chat WITHOUT
    a Flask restart.

    Best-effort: if Flask isn't running we just log and move on.
    """
    import urllib.request, urllib.error
    url = os.environ.get("FLASK_RELOAD_URL", "http://localhost:5000/api/admin/reload-index")
    try:
        req = urllib.request.Request(url, method="POST")
        with urllib.request.urlopen(req, timeout=30) as resp:
            body = resp.read().decode("utf-8", errors="replace")
            log.info(f"  Local index reload: {body[:200]}")
    except urllib.error.URLError as e:
        log.warning(f"  Local index reload skipped (Flask not reachable at {url}): {e.reason}")
    except Exception as e:
        log.warning(f"  Local index reload failed: {type(e).__name__}: {e}")


def run_scheduled(interval_minutes):
    log.info(f"Scheduled mode: syncing every {interval_minutes} minute(s). Ctrl+C to stop.")
    while True:
        try:
            run_sync()
            # After a successful sync, ping the running website so it picks
            # up the new files in its local index without needing a restart.
            trigger_local_index_reload()
        except Exception as e:
            log.error(f"Sync cycle failed: {e}")
            if "auth" in str(e).lower() or "token" in str(e).lower():
                log.error("TOKEN ERROR -- run bootstrap_onedrive.py to re-authenticate.")
        log.info(f"Next sync in {interval_minutes} minute(s)...")
        time.sleep(interval_minutes * 60)

def main():
    parser = argparse.ArgumentParser(description="OneDrive -> GCS -> Vertex AI Search sync")
    parser.add_argument("--dry-run",  action="store_true")
    parser.add_argument("--force",    action="store_true",
                        help="Re-fetch all files from OneDrive (ignores delta link)")
    parser.add_argument("--rebuild-only", action="store_true",
                        help="Skip OneDrive download. Re-build Vertex manifest from existing "
                             "GCS contents with pre-extraction + OCR. Fast (~10-30 min).")
    parser.add_argument("--schedule", type=int, metavar="MINUTES")
    args = parser.parse_args()
    if args.schedule:
        run_scheduled(args.schedule)
    else:
        run_sync(dry_run=args.dry_run, force=args.force, rebuild_only=args.rebuild_only)

if __name__ == "__main__":
    main()
